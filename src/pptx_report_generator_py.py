"""
Pure Python PPTX Report Generator
==================================
Generates premium PowerPoint reports using python-pptx.
No Node.js dependency — works on Vercel serverless.
"""

import os
import logging
from datetime import datetime

import pandas as pd
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.chart import XL_CHART_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logger = logging.getLogger(__name__)

# ── Color Palette ─────────────────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x17, 0x29) if HAS_PPTX else None
NAVY_LIGHT = RGBColor(0x1A, 0x25, 0x3C) if HAS_PPTX else None
TEAL = RGBColor(0x00, 0xB4, 0xD8) if HAS_PPTX else None
GOLD = RGBColor(0xF4, 0xA2, 0x61) if HAS_PPTX else None
EMERALD = RGBColor(0x2E, 0xCC, 0x71) if HAS_PPTX else None
ROSE = RGBColor(0xE7, 0x4C, 0x3C) if HAS_PPTX else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if HAS_PPTX else None
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0) if HAS_PPTX else None
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99) if HAS_PPTX else None
DARK_TEXT = RGBColor(0x33, 0x33, 0x33) if HAS_PPTX else None


def _safe_sum(series):
    if series is None or (hasattr(series, 'empty') and series.empty):
        return 0
    try:
        return float(series.sum())
    except Exception:
        return 0


def _safe_div(a, b, mult=1):
    try:
        if b and b != 0:
            return (a / b) * mult
    except Exception:
        pass
    return 0


def _fmt_number(n, decimals=0):
    """Format number with commas."""
    if n >= 1_000_000:
        return f"{n/1_000_000:,.{decimals}f}M"
    elif n >= 1_000:
        return f"{n/1_000:,.{decimals}f}K"
    return f"{n:,.{decimals}f}"


def _fmt_currency(n):
    """Format currency."""
    if n >= 1_000_000:
        return f"₹{n/1_000_000:,.1f}M"
    elif n >= 1_000:
        return f"₹{n/1_000:,.1f}K"
    return f"₹{n:,.0f}"


def _extract_metrics(df, platform="meta"):
    """Extract metrics from a DataFrame."""
    if df is None or df.empty:
        return {"spend": 0, "impressions": 0, "reach": 0, "clicks": 0,
                "ctr": 0, "cpm": 0, "cpc": 0, "engagements": 0}

    if platform == "meta":
        spend = _safe_sum(df["Amount spent (INR)"]) if "Amount spent (INR)" in df.columns else 0
        impressions = _safe_sum(df["Impressions"]) if "Impressions" in df.columns else 0
        reach = _safe_sum(df["Reach"]) if "Reach" in df.columns else 0
        clicks = _safe_sum(df["Clicks (all)"]) if "Clicks (all)" in df.columns else 0
        engagements = _safe_sum(df["Post engagements"]) if "Post engagements" in df.columns else 0
    else:
        spend = _safe_sum(df["Cost"]) if "Cost" in df.columns else 0
        impressions = _safe_sum(df["Impr."]) if "Impr." in df.columns else 0
        reach = impressions * 0.65
        clicks = _safe_sum(df["Clicks"]) if "Clicks" in df.columns else 0
        engagements = _safe_sum(df["TrueView views"]) if "TrueView views" in df.columns else 0

    return {
        "spend": spend,
        "impressions": impressions,
        "reach": reach,
        "clicks": clicks,
        "ctr": _safe_div(clicks, impressions, 100),
        "cpm": _safe_div(spend, impressions, 1000),
        "cpc": _safe_div(spend, clicks),
        "engagements": engagements,
    }


def _get_campaign_summary(df, name_col, spend_col, impr_col, clicks_col):
    """Group by campaign and return list of dicts."""
    if df is None or df.empty or name_col not in df.columns:
        return []

    cols = {spend_col: "sum", impr_col: "sum", clicks_col: "sum"}
    available = {k: v for k, v in cols.items() if k in df.columns}
    if not available:
        return []

    grouped = df.groupby(name_col).agg(available).reset_index()
    campaigns = []
    for _, row in grouped.iterrows():
        spend = float(row.get(spend_col, 0) or 0)
        impr = float(row.get(impr_col, 0) or 0)
        clicks = float(row.get(clicks_col, 0) or 0)
        campaigns.append({
            "name": str(row[name_col])[:40],
            "spend": spend,
            "impressions": impr,
            "clicks": clicks,
            "ctr": _safe_div(clicks, impr, 100),
            "cpm": _safe_div(spend, impr, 1000),
        })
    return sorted(campaigns, key=lambda x: x["spend"], reverse=True)[:10]


def _set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_kpi_card(slide, left, top, width, height, label, value, bg_color=NAVY_LIGHT):
    """Add a KPI card with label and value."""
    card = _add_shape(slide, left, top, width, height, fill_color=bg_color)
    card.shadow.inherit = False

    # Value
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.15),
                  width - Inches(0.3), Inches(0.5),
                  str(value), font_size=22, font_color=TEAL, bold=True,
                  alignment=PP_ALIGN.CENTER)
    # Label
    _add_text_box(slide, left + Inches(0.15), top + Inches(0.6),
                  width - Inches(0.3), Inches(0.3),
                  label, font_size=10, font_color=MEDIUM_GRAY,
                  alignment=PP_ALIGN.CENTER)


# ── Slide Generators ──────────────────────────────────────────────────

def _add_cover_slide(prs, brand_name, report_type, meta_range, yt_range):
    """Slide 1: Cover page with navy background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_bg(slide, NAVY)

    # Accent line
    _add_shape(slide, Inches(1.5), Inches(2.0), Inches(1.5), Pt(3), fill_color=TEAL)

    # Brand name
    _add_text_box(slide, Inches(1.5), Inches(2.2), Inches(7), Inches(0.8),
                  brand_name.upper(), font_size=36, font_color=WHITE, bold=True)

    # Report type
    _add_text_box(slide, Inches(1.5), Inches(3.0), Inches(7), Inches(0.5),
                  f"{report_type.upper()} CAMPAIGN PERFORMANCE REPORT",
                  font_size=16, font_color=TEAL)

    # Date ranges
    _add_text_box(slide, Inches(1.5), Inches(3.7), Inches(7), Inches(0.4),
                  f"Meta Ads: {meta_range}", font_size=11, font_color=MEDIUM_GRAY)
    _add_text_box(slide, Inches(1.5), Inches(4.0), Inches(7), Inches(0.4),
                  f"YouTube Ads: {yt_range}", font_size=11, font_color=MEDIUM_GRAY)

    # Confidential notice
    _add_text_box(slide, Inches(1.5), Inches(5.5), Inches(7), Inches(0.3),
                  "CONFIDENTIAL — PREPARED BY ADFLOW STUDIO",
                  font_size=9, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.LEFT)

    # Generated date
    _add_text_box(slide, Inches(1.5), Inches(5.8), Inches(7), Inches(0.3),
                  f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
                  font_size=9, font_color=MEDIUM_GRAY)


def _add_executive_summary(prs, total_spend, total_impr, total_clicks,
                            total_reach, total_ctr, total_cpm):
    """Slide 2: Executive Summary with KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    # Title
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "EXECUTIVE SUMMARY", font_size=24, font_color=WHITE, bold=True)
    _add_shape(slide, Inches(0.5), Inches(0.85), Inches(1.2), Pt(3), fill_color=TEAL)

    # 2x3 KPI grid
    kpis = [
        ("Total Spend", _fmt_currency(total_spend)),
        ("Impressions", _fmt_number(total_impr)),
        ("Clicks", _fmt_number(total_clicks)),
        ("Reach", _fmt_number(total_reach)),
        ("CTR", f"{total_ctr:.2f}%"),
        ("CPM", _fmt_currency(total_cpm)),
    ]

    for i, (label, value) in enumerate(kpis):
        col = i % 3
        row = i // 3
        left = Inches(0.5 + col * 3.1)
        top = Inches(1.3 + row * 1.6)
        _add_kpi_card(slide, left, top, Inches(2.8), Inches(1.2), label, value)


def _add_campaign_table_slide(prs, title, campaigns, platform_color):
    """Add a campaign performance table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    # Title
    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  title, font_size=22, font_color=WHITE, bold=True)
    _add_shape(slide, Inches(0.5), Inches(0.8), Inches(1.2), Pt(3), fill_color=platform_color)

    if not campaigns:
        _add_text_box(slide, Inches(1), Inches(2.5), Inches(8), Inches(0.5),
                      "No campaign data available", font_size=14, font_color=MEDIUM_GRAY,
                      alignment=PP_ALIGN.CENTER)
        return

    # Table
    headers = ["Campaign", "Spend", "Impressions", "Clicks", "CTR", "CPM"]
    rows_data = campaigns[:8]

    table_shape = slide.shapes.add_table(
        len(rows_data) + 1, len(headers),
        Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.4 * (len(rows_data) + 1))
    )
    table = table_shape.table

    # Style header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = platform_color

    # Style rows
    for i, camp in enumerate(rows_data):
        row_idx = i + 1
        values = [
            camp["name"][:35],
            _fmt_currency(camp["spend"]),
            _fmt_number(camp["impressions"]),
            _fmt_number(camp["clicks"]),
            f"{camp['ctr']:.2f}%",
            _fmt_currency(camp["cpm"]),
        ]
        bg = NAVY_LIGHT if i % 2 == 0 else NAVY
        for j, val in enumerate(values):
            cell = table.cell(row_idx, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = WHITE
                paragraph.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Set column widths
    col_widths = [Inches(3.0), Inches(1.2), Inches(1.5), Inches(1.2), Inches(1.0), Inches(1.5)]
    for j, w in enumerate(col_widths):
        table.columns[j].width = w


def _add_platform_comparison(prs, meta_metrics, yt_metrics):
    """Slide: Platform comparison side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "PLATFORM COMPARISON", font_size=22, font_color=WHITE, bold=True)
    _add_shape(slide, Inches(0.5), Inches(0.8), Inches(1.2), Pt(3), fill_color=TEAL)

    # Meta column
    _add_shape(slide, Inches(0.5), Inches(1.2), Inches(4.2), Inches(4.5), fill_color=NAVY_LIGHT)
    _add_text_box(slide, Inches(0.7), Inches(1.3), Inches(3.8), Inches(0.4),
                  "META ADS", font_size=16, font_color=RGBColor(0x1B, 0x77, 0xF2), bold=True,
                  alignment=PP_ALIGN.CENTER)

    meta_items = [
        ("Spend", _fmt_currency(meta_metrics["spend"])),
        ("Impressions", _fmt_number(meta_metrics["impressions"])),
        ("Clicks", _fmt_number(meta_metrics["clicks"])),
        ("CTR", f"{meta_metrics['ctr']:.2f}%"),
        ("CPM", _fmt_currency(meta_metrics["cpm"])),
        ("CPC", _fmt_currency(meta_metrics["cpc"])),
    ]
    for i, (label, val) in enumerate(meta_items):
        y = Inches(1.9 + i * 0.6)
        _add_text_box(slide, Inches(0.8), y, Inches(1.8), Inches(0.3),
                      label, font_size=10, font_color=MEDIUM_GRAY)
        _add_text_box(slide, Inches(2.8), y, Inches(1.8), Inches(0.3),
                      val, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

    # YouTube column
    _add_shape(slide, Inches(5.3), Inches(1.2), Inches(4.2), Inches(4.5), fill_color=NAVY_LIGHT)
    _add_text_box(slide, Inches(5.5), Inches(1.3), Inches(3.8), Inches(0.4),
                  "YOUTUBE ADS", font_size=16, font_color=RGBColor(0xFF, 0x00, 0x00), bold=True,
                  alignment=PP_ALIGN.CENTER)

    yt_items = [
        ("Spend", _fmt_currency(yt_metrics["spend"])),
        ("Impressions", _fmt_number(yt_metrics["impressions"])),
        ("Clicks", _fmt_number(yt_metrics["clicks"])),
        ("CTR", f"{yt_metrics['ctr']:.2f}%"),
        ("CPM", _fmt_currency(yt_metrics["cpm"])),
        ("CPC", _fmt_currency(yt_metrics["cpc"])),
    ]
    for i, (label, val) in enumerate(yt_items):
        y = Inches(1.9 + i * 0.6)
        _add_text_box(slide, Inches(5.6), y, Inches(1.8), Inches(0.3),
                      label, font_size=10, font_color=MEDIUM_GRAY)
        _add_text_box(slide, Inches(7.6), y, Inches(1.8), Inches(0.3),
                      val, font_size=12, font_color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)


def _add_budget_pacing(prs, meta_spend, meta_budget, yt_spend, yt_budget):
    """Slide: Budget pacing with progress indicators."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "BUDGET PACING", font_size=22, font_color=WHITE, bold=True)
    _add_shape(slide, Inches(0.5), Inches(0.8), Inches(1.2), Pt(3), fill_color=GOLD)

    # Meta budget
    meta_pct = min(_safe_div(meta_spend, meta_budget, 100), 100) if meta_budget else 0
    _add_shape(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(2.0), fill_color=NAVY_LIGHT)
    _add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4), Inches(0.4),
                  "Meta Ads Budget", font_size=14, font_color=WHITE, bold=True)
    _add_text_box(slide, Inches(5), Inches(1.5), Inches(4.2), Inches(0.4),
                  f"{_fmt_currency(meta_spend)} / {_fmt_currency(meta_budget)} ({meta_pct:.0f}%)",
                  font_size=12, font_color=TEAL, alignment=PP_ALIGN.RIGHT)

    # Progress bar background
    _add_shape(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.3), fill_color=NAVY)
    # Progress bar fill
    bar_width = max(Inches(0.1), Inches(8.4 * meta_pct / 100))
    bar_color = EMERALD if meta_pct < 90 else (GOLD if meta_pct < 100 else ROSE)
    _add_shape(slide, Inches(0.8), Inches(2.2), bar_width, Inches(0.3), fill_color=bar_color)

    # YouTube budget
    yt_pct = min(_safe_div(yt_spend, yt_budget, 100), 100) if yt_budget else 0
    _add_shape(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(2.0), fill_color=NAVY_LIGHT)
    _add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4), Inches(0.4),
                  "YouTube Ads Budget", font_size=14, font_color=WHITE, bold=True)
    _add_text_box(slide, Inches(5), Inches(3.8), Inches(4.2), Inches(0.4),
                  f"{_fmt_currency(yt_spend)} / {_fmt_currency(yt_budget)} ({yt_pct:.0f}%)",
                  font_size=12, font_color=TEAL, alignment=PP_ALIGN.RIGHT)

    _add_shape(slide, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.3), fill_color=NAVY)
    bar_width_yt = max(Inches(0.1), Inches(8.4 * yt_pct / 100))
    bar_color_yt = EMERALD if yt_pct < 90 else (GOLD if yt_pct < 100 else ROSE)
    _add_shape(slide, Inches(0.8), Inches(4.5), bar_width_yt, Inches(0.3), fill_color=bar_color_yt)


def _add_insights_slide(prs, meta_metrics, yt_metrics, total_spend, total_ctr):
    """Slide: Key insights and recommendations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                  "KEY INSIGHTS & RECOMMENDATIONS", font_size=22, font_color=WHITE, bold=True)
    _add_shape(slide, Inches(0.5), Inches(0.8), Inches(1.2), Pt(3), fill_color=EMERALD)

    insights = []

    # Generate dynamic insights
    if total_ctr > 2:
        insights.append(("High CTR Performance", f"Combined CTR of {total_ctr:.2f}% exceeds industry average of 1.5%.", EMERALD))
    elif total_ctr > 1:
        insights.append(("Moderate CTR", f"CTR of {total_ctr:.2f}% is near industry average. Consider A/B testing creatives.", GOLD))
    else:
        insights.append(("Low CTR Alert", f"CTR of {total_ctr:.2f}% is below benchmark. Review targeting and ad copy.", ROSE))

    if meta_metrics["spend"] > 0 and yt_metrics["spend"] > 0:
        meta_share = meta_metrics["spend"] / total_spend * 100
        insights.append(("Budget Split", f"Meta receives {meta_share:.0f}% of total spend. Consider rebalancing based on ROAS.", TEAL))

    if meta_metrics["cpc"] > 0:
        insights.append(("Meta CPC", f"Average CPC is {_fmt_currency(meta_metrics['cpc'])}. Optimize bids for high-performing segments.", TEAL))

    if yt_metrics["impressions"] > 0:
        insights.append(("YouTube Reach", f"{_fmt_number(yt_metrics['impressions'])} impressions delivered. Focus on view-through conversions.", TEAL))

    insights.append(("Next Steps", "Schedule weekly performance reviews and set up automated alerts for KPI thresholds.", MEDIUM_GRAY))

    for i, (title, desc, color) in enumerate(insights[:5]):
        y = Inches(1.2 + i * 1.0)
        _add_shape(slide, Inches(0.5), y, Inches(9), Inches(0.8), fill_color=NAVY_LIGHT)
        # Color accent
        _add_shape(slide, Inches(0.5), y, Pt(4), Inches(0.8), fill_color=color)
        _add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(8.5), Inches(0.3),
                      title, font_size=12, font_color=color, bold=True)
        _add_text_box(slide, Inches(0.8), y + Inches(0.35), Inches(8.5), Inches(0.35),
                      desc, font_size=10, font_color=LIGHT_GRAY)


def _add_closing_slide(prs, brand_name):
    """Final slide: Thank you / closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_shape(slide, Inches(3), Inches(2.0), Inches(4), Pt(3), fill_color=TEAL)

    _add_text_box(slide, Inches(1), Inches(2.3), Inches(8), Inches(0.8),
                  "THANK YOU", font_size=40, font_color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.5),
                  f"Report prepared for {brand_name}",
                  font_size=14, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.4),
                  "Powered by AdFlow Studio — Campaign Intelligence Platform",
                  font_size=11, font_color=TEAL, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.3),
                  f"Generated {datetime.now().strftime('%B %d, %Y')} | Confidential",
                  font_size=9, font_color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ── Main Generator ────────────────────────────────────────────────────

def generate_pptx_report(meta_data, google_data, brand_config, report_type="full", output_path=None):
    """Generate a premium PPTX report using python-pptx.

    Args:
        meta_data: dict with raw_data, campaign_data, etc.
        google_data: dict with raw_data, campaign_data, etc.
        brand_config: brand configuration dict
        report_type: 'daily', 'weekly', 'monthly', 'full'
        output_path: where to save the .pptx file

    Returns:
        path to generated PPTX file, or None on failure
    """
    if not HAS_PPTX:
        logger.error("python-pptx not installed")
        return None

    if output_path is None:
        output_path = f"/tmp/report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    try:
        # Extract data
        meta_raw = meta_data.get("raw_data", pd.DataFrame()) if meta_data else pd.DataFrame()
        yt_raw = google_data.get("raw_data", pd.DataFrame()) if google_data else pd.DataFrame()

        meta_metrics = _extract_metrics(meta_raw, "meta")
        yt_metrics = _extract_metrics(yt_raw, "youtube")

        total_spend = meta_metrics["spend"] + yt_metrics["spend"]
        total_impr = meta_metrics["impressions"] + yt_metrics["impressions"]
        total_clicks = meta_metrics["clicks"] + yt_metrics["clicks"]
        total_reach = meta_metrics["reach"] + yt_metrics["reach"]
        total_ctr = _safe_div(total_clicks, total_impr, 100)
        total_cpm = _safe_div(total_spend, total_impr, 1000)

        brand_name = brand_config.get("name", "Campaign Report") if isinstance(brand_config, dict) else "Campaign Report"
        meta_cfg = brand_config.get("meta", {}) if isinstance(brand_config, dict) else {}
        yt_cfg = brand_config.get("youtube", {}) if isinstance(brand_config, dict) else {}

        meta_range = f"{meta_cfg.get('start_date', 'N/A')} to {meta_cfg.get('end_date', 'N/A')}"
        yt_range = f"{yt_cfg.get('start_date', 'N/A')} to {yt_cfg.get('end_date', 'N/A')}"

        meta_budget = meta_cfg.get("budget", 0) or 0
        yt_budget = yt_cfg.get("budget", 0) or 0

        # Create presentation (16:9 widescreen)
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Build slides
        _add_cover_slide(prs, brand_name, report_type, meta_range, yt_range)
        _add_executive_summary(prs, total_spend, total_impr, total_clicks,
                               total_reach, total_ctr, total_cpm)

        # Meta campaigns
        meta_campaigns = _get_campaign_summary(meta_raw, "Campaign name",
                                                "Amount spent (INR)", "Impressions", "Clicks (all)")
        _add_campaign_table_slide(prs, "META ADS PERFORMANCE", meta_campaigns,
                                  RGBColor(0x1B, 0x77, 0xF2))

        # YouTube campaigns
        yt_campaigns = _get_campaign_summary(yt_raw, "Campaign", "Cost", "Impr.", "Clicks")
        _add_campaign_table_slide(prs, "YOUTUBE ADS PERFORMANCE", yt_campaigns,
                                  RGBColor(0xFF, 0x00, 0x00))

        # Platform comparison
        _add_platform_comparison(prs, meta_metrics, yt_metrics)

        # Budget pacing
        _add_budget_pacing(prs, meta_metrics["spend"], meta_budget,
                          yt_metrics["spend"], yt_budget)

        # Insights
        _add_insights_slide(prs, meta_metrics, yt_metrics, total_spend, total_ctr)

        # Closing
        _add_closing_slide(prs, brand_name)

        # Save
        prs.save(output_path)
        logger.info(f"PPTX report generated: {output_path}")
        return output_path

    except Exception as e:
        logger.error(f"PPTX generation error: {e}")
        return None
